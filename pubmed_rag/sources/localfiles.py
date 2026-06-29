"""Read a local document corpus (a directory tree) into plain-text Documents.

Supported: .pdf (pdfplumber), .docx (python-docx), .pptx (python-pptx), and
.txt/.md (read directly). The office-format libraries are lazy-imported per type
so the core package never hard-depends on them. Files with no extractable text
(e.g. image-only/scanned PDFs) are skipped via the on_skip callback -- nothing
is silently dropped.

Retrieval/embedding is NOT done here; this layer only turns files into text.
"""
import re
from dataclasses import dataclass, field
from pathlib import Path

TEXT_EXT = {".txt", ".md"}
SUPPORTED = {".pdf", ".docx", ".pptx", ".xlsx", *TEXT_EXT}


@dataclass
class Document:
    doc_id: str          # stable id = path relative to the corpus root
    title: str
    text: str
    source_path: str
    category: str        # top-level subfolder (audience/topic facet), "" if at root
    metadata: dict = field(default_factory=dict)


def iter_files(root, on_skip=None):
    """Yield (rel_path, Path, category) for every SUPPORTED file under `root`,
    WITHOUT extracting. Cheap -- used to hash files for incremental ingest and to
    count work for a progress bar before the slow extraction step."""
    root = Path(root)
    if not root.exists():
        raise SystemExit(f"Corpus directory not found: {root}")
    for path in sorted(p for p in root.rglob("*") if p.is_file()):
        rel = str(path.relative_to(root))
        ext = path.suffix.lower()
        if ext not in SUPPORTED:
            if on_skip:
                on_skip(rel, f"unsupported type {ext or '(none)'}")
            continue
        category = rel.split("/")[0] if "/" in rel.replace("\\", "/") else ""
        yield rel, path, category


def extract_document(rel, path, category, on_skip=None):
    """Extract one file to a Document, or None (empty / malformed -> on_skip)."""
    ext = path.suffix.lower()
    try:
        text = _clean(_extract(path, ext))
    except Exception as e:  # a malformed file shouldn't abort the whole run
        if on_skip:
            on_skip(rel, f"extract error: {e}")
        return None
    if not text:
        if on_skip:
            on_skip(rel, "no extractable text (image-only?)")
        return None
    return Document(
        doc_id=rel,
        title=path.stem,
        text=text,
        source_path=str(path),
        category=category,
        metadata={"ext": ext, "chars": len(text)},
    )


def iter_documents(root, on_skip=None):
    """Yield a Document per supported file with extractable text (discovery + extraction)."""
    for rel, path, category in iter_files(root, on_skip):
        doc = extract_document(rel, path, category, on_skip)
        if doc:
            yield doc


def _extract(path: Path, ext: str) -> str:
    if ext in TEXT_EXT:
        return path.read_text(encoding="utf-8", errors="ignore")
    if ext == ".pdf":
        return _pdf(path)
    if ext == ".docx":
        return _docx(path)
    if ext == ".pptx":
        return _pptx(path)
    if ext == ".xlsx":
        return _xlsx(path)
    return ""


def _pdf(path: Path) -> str:
    import logging

    logging.getLogger("pdfminer").setLevel(logging.ERROR)  # silence FontBBox noise
    import pdfplumber

    with pdfplumber.open(path) as pdf:
        return "\n\n".join((page.extract_text() or "") for page in pdf.pages)


def _docx(path: Path) -> str:
    import docx

    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:  # flatten simple tables row by row
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            s.text.strip()
            for s in slide.shapes
            if s.has_text_frame and s.text.strip()
        ]
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        parts = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
        return "\n\n".join(parts)
    finally:
        wb.close()


_WS = re.compile(r"[ \t]+")
_NL = re.compile(r"\n{3,}")


def _clean(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _WS.sub(" ", text)
    text = _NL.sub("\n\n", text)
    return "\n".join(line.strip() for line in text.split("\n")).strip()
